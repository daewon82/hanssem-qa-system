from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HANSSEM_BLUE  = RGBColor(0x00, 0x5B, 0xAC)
HANSSEM_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
LIGHT_BLUE    = RGBColor(0xE8, 0xF4, 0xFF)
LIGHT_GRAY    = RGBColor(0xF5, 0xF5, 0xF5)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT     = RGBColor(0x2C, 0x2C, 0x2C)
MID_GRAY      = RGBColor(0x6B, 0x6B, 0x6B)
GREEN         = RGBColor(0x00, 0xB3, 0x5A)
RED           = RGBColor(0xE5, 0x39, 0x35)
YELLOW        = RGBColor(0xFF, 0xC1, 0x07)
LIGHT_GREEN   = RGBColor(0xE8, 0xF8, 0xEF)
LIGHT_RED     = RGBColor(0xFD, 0xED, 0xEC)
LIGHT_YELLOW  = RGBColor(0xFF, 0xF8, 0xE1)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def rect(slide, x, y, w, h, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    return s

def tb(slide, text, x, y, w, h, sz=11, bold=False, color=DARK_TEXT,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "맑은 고딕"
    return box

def ml(slide, lines, x, y, w, h, sz=10, bold=False, color=DARK_TEXT,
        align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "맑은 고딕"
    return box

# ══════════════════════════════════════════════════════════════
# 슬라이드: 사내망 이전 계획
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)

# 상단 헤더 바
rect(slide, 0, 0, 13.33, 1.1, fill=HANSSEM_BLUE)
rect(slide, 0.25, 0.15, 0.75, 0.75, fill=ACCENT_ORANGE)
tb(slide, "09. 향후 계획 — 대시보드 사내망 이전", 1.2, 0.12, 10.5, 0.55,
   sz=24, bold=True, color=WHITE)
tb(slide, "Dashboard Migration: GitHub Pages → Internal Network", 1.2, 0.65, 10.5, 0.38,
   sz=11, color=RGBColor(0xCC, 0xE4, 0xFF))

# ── 현재 구조 (AS-IS) ────────────────────────────────────────
rect(slide, 0.3, 1.2, 5.8, 4.8, fill=WHITE, line=HANSSEM_BLUE, lw=0.5)
rect(slide, 0.3, 1.2, 5.8, 0.42, fill=HANSSEM_BLUE)
tb(slide, "현재 구조 (AS-IS)", 0.45, 1.23, 5.5, 0.36,
   sz=12, bold=True, color=WHITE)

as_is = [
    "GitHub Actions",
    "  └─ 테스트 실행 (PC + MW)",
    "  └─ publishResults() → gh-pages에 직접 기록",
    "  └─ updateProgress() → gh-pages에 직접 기록",
    "  └─ git push → main 브랜치",
    "  └─ peaceiris/actions-gh-pages@v3",
    "       → public/ 폴더를 gh-pages 브랜치로 배포",
    "",
    "대시보드 (index.html)",
    "  └─ https://daewon82.github.io/hanssem-qa-system/",
    "  └─ progress.json URL: gh-pages 하드코딩",
    "  └─ GitHub API 폴링 (워크플로우 상태 확인)",
    "  └─ PAT 토큰: index.html에 내장",
]
ml(slide, as_is, 0.45, 1.7, 5.55, 4.1, sz=10, color=DARK_TEXT)

# 화살표
tb(slide, "→", 6.25, 3.25, 0.6, 0.6, sz=28, bold=True,
   color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
tb(slide, "이전", 6.2, 3.8, 0.7, 0.3, sz=9, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

# ── 이전 구조 (TO-BE) ────────────────────────────────────────
rect(slide, 7.1, 1.2, 5.9, 4.8, fill=WHITE, line=GREEN, lw=0.5)
rect(slide, 7.1, 1.2, 5.9, 0.42, fill=GREEN)
tb(slide, "이전 구조 (TO-BE)", 7.25, 1.23, 5.6, 0.36,
   sz=12, bold=True, color=WHITE)

to_be = [
    "GitHub Actions  ← 변경 없음",
    "  └─ 테스트 실행 (PC + MW)  ← 변경 없음",
    "  └─ publishResults()  ← 내부 서버 URL로 수정 필요",
    "  └─ updateProgress()  ← 내부 서버 URL로 수정 필요",
    "  └─ git push → main 브랜치  ← 변경 없음",
    "  └─ 배포 방식 변경",
    "       기존: peaceiris/actions-gh-pages@v3",
    "       변경: scp / rsync → 내부 웹 서버",
    "",
    "대시보드 (index.html)",
    "  └─ https://내부망주소/qa-system/  ← URL 변경",
    "  └─ progress.json URL  ← 내부 서버 URL로 수정",
    "  └─ GitHub API 폴링  ← 내부망→GitHub 가능 여부 확인",
    "  └─ PAT 토큰  ← 내부망→GitHub 접근 필요 여부 확인",
]
ml(slide, to_be, 7.25, 1.7, 5.65, 4.1, sz=10, color=DARK_TEXT)

# ── 하단 변경 요약 ────────────────────────────────────────────
rect(slide, 0.3, 6.15, 12.73, 1.2, fill=WHITE, line=HANSSEM_BLUE, lw=0.3)
rect(slide, 0.3, 6.15, 12.73, 0.38, fill=HANSSEM_DARK)
tb(slide, "변경 범위 요약", 0.45, 6.18, 12.4, 0.32, sz=11, bold=True, color=WHITE)

items = [
    ("변경 없음 ✅", LIGHT_GREEN, GREEN,
     "테스트 코드 전체 / GitHub Actions 워크플로우 / 결과 JSON 구조"),
    ("수정 필요 ⚠️", LIGHT_YELLOW, YELLOW,
     "utils.ts — publishResults·updateProgress의 gh-pages URL\nindex.html — progress.json URL, GitHub API 폴링 URL"),
    ("변경 필요 🔄", LIGHT_RED, RED,
     "배포 방식 — actions-gh-pages → 내부 서버 전송(scp/rsync)\n내부망→GitHub API 접근 가능 여부에 따라 PAT 정책 재검토"),
]
for i, (label, bg, border, desc) in enumerate(items):
    x = 0.4 + i * 4.25
    rect(slide, x, 6.6, 4.1, 0.65, fill=bg, line=border, lw=0.4)
    tb(slide, label, x+0.08, 6.62, 1.5, 0.28, sz=9, bold=True, color=border)
    tb(slide, desc, x+0.08, 6.88, 3.95, 0.35, sz=8.5, color=DARK_TEXT)

prs.save("/Users/dw/hanssem-qa-system/사내망_이전계획.pptx")
print("✅ PPT 저장 완료: 사내망_이전계획.pptx")
