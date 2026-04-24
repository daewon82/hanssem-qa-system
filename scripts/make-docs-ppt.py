"""한샘몰 서비스 품질 관리 대시보드 — PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── 색상 팔레트 ───
PRIMARY     = RGBColor(0x1A, 0x4F, 0x9C)  # 한샘 블루
ACCENT      = RGBColor(0xE6, 0x51, 0x00)  # 주황
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
RED         = RGBColor(0xC6, 0x28, 0x28)
AMBER       = RGBColor(0xE6, 0x51, 0x00)
PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)
TEAL        = RGBColor(0x00, 0x69, 0x5C)
TEXT_MAIN   = RGBColor(0x1A, 0x1E, 0x2D)
TEXT_SUB    = RGBColor(0x4B, 0x55, 0x68)
TEXT_HINT   = RGBColor(0x9C, 0xA3, 0xAF)
BG_CARD     = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT    = RGBColor(0xF4, 0xF6, 0xFB)
BORDER      = RGBColor(0xE5, 0xE7, 0xEB)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

W = prs.slide_width
H = prs.slide_height


# ─── 유틸 ───
def add_rect(slide, x, y, w, h, fill=BG_CARD, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=TEXT_MAIN,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, name="Pretendard"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    return tb


def add_bullets(slide, x, y, w, h, lines, size=14, color=TEXT_SUB, name="Pretendard"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, (prefix, text) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"{prefix} {text}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = name
    return tb


def header_bar(slide, title, subtitle=None):
    # 상단 색상 바
    add_rect(slide, 0, 0, W, Inches(0.28), fill=PRIMARY)
    # 타이틀
    add_text(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.7),
             title, size=28, bold=True, color=TEXT_MAIN)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.35),
                 subtitle, size=13, color=TEXT_HINT)


def footer(slide, page):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "한샘몰 QA 시스템 · 자비스 PM", size=9, color=TEXT_HINT)
    add_text(slide, Inches(12), Inches(7.1), Inches(1.3), Inches(0.3),
             f"{page:02d}", size=9, color=TEXT_HINT, align=PP_ALIGN.RIGHT)


page_no = [0]
def next_page():
    page_no[0] += 1
    return page_no[0]


# ═══════════════════════════════════════════════════════════════
# 슬라이드 1: 타이틀
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
# 한샘 블루 강조 밴드
add_rect(s, 0, Inches(2.8), W, Inches(2.0), fill=PRIMARY)
add_text(s, Inches(0.5), Inches(3.1), Inches(12.3), Inches(0.6),
         "HANSSEM MALL", size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(1.0),
         "서비스 품질 관리 대시보드", size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
         "전체 구조 · 테스트 기능 · 고도화 전략", size=18, color=RGBColor(0xFF, 0xFF, 0xFF, 0xE0) if False else RGBColor(0xE0, 0xE5, 0xF0))
add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
         "2026-04-23 · v1.0 · JARVIS PM", size=11, color=TEXT_SUB)


# ═══════════════════════════════════════════════════════════════
# 슬라이드 2: 프로젝트 개요
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
header_bar(s, "01 · 프로젝트 개요", "한샘몰 운영환경 자동 QA 플랫폼")

add_text(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
         "🎯 목적", size=18, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.7), Inches(2.3), Inches(12), Inches(1.5), [
    ("▸", "한샘몰(store / m.store / mall) 전 구간 가용성·기능 회귀 감지"),
    ("▸", "매일 자동 실행으로 서비스 품질 변화 추적"),
    ("▸", "실패 원인의 AI 기반 자동 분석으로 개발팀 대응 시간 단축"),
], size=14)

add_text(s, Inches(0.5), Inches(4.0), Inches(12), Inches(0.4),
         "🛠 핵심 스택", size=18, bold=True, color=PRIMARY)

# 스택 카드
items = [
    ("Playwright", "+ TypeScript", PRIMARY),
    ("GitHub Actions", "Ubuntu 러너", ACCENT),
    ("GitHub Pages", "정적 대시보드", GREEN),
    ("Jandi Webhook", "팀 알림", PURPLE),
    ("Claude Sonnet 4.6", "AI 실패 분석", TEAL),
]
card_w = Inches(2.4); card_h = Inches(1.4); gap = Inches(0.1)
start_x = Inches(0.5)
for i, (name, desc, color) in enumerate(items):
    x = start_x + i * (card_w + gap)
    add_rect(s, x, Inches(4.6), card_w, card_h, fill=BG_CARD, line=BORDER)
    add_rect(s, x, Inches(4.6), card_w, Inches(0.08), fill=color)
    add_text(s, x + Inches(0.15), Inches(4.8), card_w - Inches(0.3), Inches(0.4),
             name, size=15, bold=True, color=TEXT_MAIN)
    add_text(s, x + Inches(0.15), Inches(5.3), card_w - Inches(0.3), Inches(0.7),
             desc, size=12, color=TEXT_SUB)

add_text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.3),
         "📍 대시보드: daewon82.github.io/hanssem-qa-system", size=11, color=TEXT_HINT)
footer(s, next_page())


# ═══════════════════════════════════════════════════════════════
# 슬라이드 3: 시스템 아키텍처
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
header_bar(s, "02 · 시스템 아키텍처", "매일 7시 자동 실행 → 5단계 테스트 → 대시보드")

# 흐름도 그리기
# 블록 1: Daily Trigger
add_rect(s, Inches(0.5), Inches(2.0), Inches(2.3), Inches(1.0), fill=PRIMARY)
add_text(s, Inches(0.5), Inches(2.0), Inches(2.3), Inches(0.4),
         "Daily Trigger", size=13, bold=True, color=BG_CARD, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.5), Inches(2.4), Inches(2.3), Inches(0.35),
         "매일 KST 07:00", size=11, color=RGBColor(0xE0, 0xE5, 0xF0), align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(2.7), Inches(2.3), Inches(0.3),
         "(cron: 22 UTC)", size=9, color=RGBColor(0xE0, 0xE5, 0xF0), align=PP_ALIGN.CENTER)

# 화살표 1
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.9), Inches(2.3), Inches(0.5), Inches(0.5))
arrow.fill.solid(); arrow.fill.fore_color.rgb = TEXT_SUB
arrow.line.fill.background()

# 블록 2: GitHub Actions
add_rect(s, Inches(3.5), Inches(1.8), Inches(5.8), Inches(3.8), fill=BG_CARD, line=BORDER)
add_text(s, Inches(3.5), Inches(1.85), Inches(5.8), Inches(0.4),
         "GitHub Actions (playwright.yml)", size=13, bold=True, color=PRIMARY,
         align=PP_ALIGN.CENTER)

# 5개 테스트 카드
tests = [
    ("1", "Crawling (PC → MW)", "100 URL × 2", PRIMARY),
    ("2", "Random (PC → MW)", "200 URL × 2 · 누적 이력 기반", PURPLE),
    ("3", "AutoE2E Public PC", "Desktop Chrome · 52 tests", ACCENT),
    ("4", "AutoE2E Authed", "Desktop + 로그인 · 31 tests", AMBER),
    ("5", "AutoE2E Public Mobile", "Pixel 5 · 52 tests", GREEN),
]
for i, (num, title, desc, color) in enumerate(tests):
    y = Inches(2.35 + i * 0.6)
    add_rect(s, Inches(3.7), y, Inches(5.4), Inches(0.52), fill=BG_LIGHT)
    # 번호 원
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.8), y + Inches(0.08), Inches(0.36), Inches(0.36))
    circle.fill.solid(); circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(s, Inches(3.8), y + Inches(0.08), Inches(0.36), Inches(0.36),
             num, size=12, bold=True, color=BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(4.3), y + Inches(0.04), Inches(3.4), Inches(0.25),
             title, size=12, bold=True, color=TEXT_MAIN)
    add_text(s, Inches(4.3), y + Inches(0.26), Inches(4.7), Inches(0.25),
             desc, size=10, color=TEXT_SUB)

# 화살표 2
arrow2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.4), Inches(2.3), Inches(0.5), Inches(0.5))
arrow2.fill.solid(); arrow2.fill.fore_color.rgb = TEXT_SUB
arrow2.line.fill.background()

# 블록 3: 결과
add_rect(s, Inches(10.0), Inches(1.8), Inches(2.9), Inches(3.8), fill=BG_CARD, line=BORDER)
add_text(s, Inches(10.0), Inches(1.9), Inches(2.9), Inches(0.4),
         "산출물", size=13, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

outs = [
    ("📊", "대시보드", "gh-pages"),
    ("🤖", "AI 실패 분석", "Claude 자동 생성"),
    ("📈", "커버리지 통계", "누적 URL·통과율"),
    ("💬", "통합 잔디 알림", "1회 전송"),
]
for i, (icon, title, desc) in enumerate(outs):
    y = Inches(2.4 + i * 0.78)
    add_text(s, Inches(10.15), y, Inches(0.4), Inches(0.4),
             icon, size=18)
    add_text(s, Inches(10.7), y, Inches(2.2), Inches(0.3),
             title, size=12, bold=True, color=TEXT_MAIN)
    add_text(s, Inches(10.7), y + Inches(0.28), Inches(2.2), Inches(0.3),
             desc, size=10, color=TEXT_SUB)

footer(s, next_page())


# ═══════════════════════════════════════════════════════════════
# 슬라이드 4: Playwright 프로젝트 구성
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
header_bar(s, "03 · Playwright 프로젝트 구성", "5개 프로젝트 순차 실행 (workers: 1)")

headers = ["Project", "대상", "실행 환경", "재시도"]
rows = [
    ("Crawling",             "PC+MW 크롤링",      "Desktop Chrome · 동적 context",  "0"),
    ("Random",               "PC+MW 랜덤 샘플",   "Desktop Chrome · 동적 context",  "0"),
    ("AutoE2E_Public_PC",    "비로그인 기능 (PC)", "Desktop Chrome",                  "1"),
    ("AutoE2E_Authed",       "로그인 필요 (PC)",  "Desktop Chrome + storageState",   "1"),
    ("AutoE2E_Public_Mobile","비로그인 기능 (MW)", "Pixel 5",                        "1"),
]

col_widths = [Inches(2.6), Inches(2.8), Inches(4.5), Inches(1.4)]
table_x = Inches(0.8)
row_h = Inches(0.6)

# 헤더
x = table_x
for i, (h, w) in enumerate(zip(headers, col_widths)):
    add_rect(s, x, Inches(2.0), w, Inches(0.5), fill=PRIMARY)
    add_text(s, x + Inches(0.15), Inches(2.0), w - Inches(0.3), Inches(0.5),
             h, size=12, bold=True, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)
    x += w

# 데이터
for r_i, row in enumerate(rows):
    x = table_x
    y = Inches(2.5 + r_i * 0.6)
    bg = BG_CARD if r_i % 2 == 0 else BG_LIGHT
    for c_i, (val, w) in enumerate(zip(row, col_widths)):
        add_rect(s, x, y, w, row_h, fill=bg, line=BORDER)
        add_text(s, x + Inches(0.15), y, w - Inches(0.3), row_h,
                 val, size=11, color=TEXT_MAIN if c_i == 0 else TEXT_SUB,
                 bold=(c_i == 0), anchor=MSO_ANCHOR.MIDDLE)
        x += w

add_text(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.4),
         "⚙ workers: 1 고정 (한샘몰 rate-limit 회피 · 순차 실행)",
         size=12, color=TEXT_HINT)
add_text(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.4),
         "📦 globalSetup: 1회 로그인 → .auth/user.json 저장 → AutoE2E_Authed 재사용",
         size=12, color=TEXT_HINT)

footer(s, next_page())


# ═══════════════════════════════════════════════════════════════
# 슬라이드 5: Crawling Test
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
header_bar(s, "04 · Crawling Test",
           "한샘몰 전체 사이트의 HTTP 가용성 전수 점검 (회귀 감지)")

# 왼쪽: 동작 방식
add_rect(s, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.8), fill=BG_CARD, line=BORDER)
add_rect(s, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.5), fill=PRIMARY)
add_text(s, Inches(0.7), Inches(1.8), Inches(5.6), Inches(0.5),
         "🔎 동작 방식", size=14, bold=True, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)

steps = [
    "① 홈(store / m.store)에서 시작",
    "② <a href> 전체 수집 (최대 1500개 풀)",
    "③ 필터: 도메인, 확장자, 경로, 키워드",
    "④ 수집 순서대로 순회 (결정적)",
    "⑤ HTTP 200 + Body 404 체크",
    "⑥ 500개 URL 완료 시 종료",
]
for i, text in enumerate(steps):
    add_text(s, Inches(0.8), Inches(2.5 + i * 0.6), Inches(5.5), Inches(0.5),
             text, size=13, color=TEXT_MAIN)

# 오른쪽: 판정 기준 + 산출물
add_rect(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(2.2), fill=BG_CARD, line=BORDER)
add_text(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(0.4),
         "📏 판정 기준", size=14, bold=True, color=PRIMARY)
criteria = [
    ("✅", "PASS", "HTTP 200 + Body 정상", GREEN),
    ("❌", "FAIL", "Timeout / 4xx,5xx / Body 404", RED),
    ("🔄", "재시도", "5xx: 10초 × 2회", AMBER),
]
for i, (icon, label, desc, color) in enumerate(criteria):
    y = Inches(2.4 + i * 0.45)
    add_text(s, Inches(7.0), y, Inches(0.3), Inches(0.4), icon, size=14)
    add_text(s, Inches(7.4), y, Inches(1.1), Inches(0.4),
             label, size=12, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.5), y, Inches(4.2), Inches(0.4),
             desc, size=11, color=TEXT_SUB, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(6.8), Inches(4.2), Inches(6.0), Inches(2.4), fill=BG_CARD, line=BORDER)
add_text(s, Inches(7.0), Inches(4.3), Inches(5.6), Inches(0.4),
         "📁 산출물", size=14, bold=True, color=PRIMARY)
outs = [
    "• public/pc_500.json · mw_500.json (전체 결과)",
    "• public/pc_url_pool.json (Random 테스트용 풀)",
    "• public/pc_tested_urls.json (누적 이력)",
]
for i, text in enumerate(outs):
    add_text(s, Inches(7.0), Inches(4.8 + i * 0.42), Inches(5.7), Inches(0.4),
             text, size=11, color=TEXT_MAIN)

add_text(s, Inches(7.0), Inches(6.2), Inches(5.7), Inches(0.4),
         "💡 매일 같은 URL을 체크 → 회귀 감지에 강함",
         size=11, color=TEXT_HINT, bold=True)

footer(s, next_page())


# ═══════════════════════════════════════════════════════════════
# 슬라이드 6: Random Test
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
header_bar(s, "05 · Random Test",
           "크롤링이 못 다루는 URL의 잠재적 이슈 발견 (Discovery)")

# 왼쪽: 동작 방식
add_rect(s, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.8), fill=BG_CARD, line=BORDER)
add_rect(s, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.5), fill=PURPLE)
add_text(s, Inches(0.7), Inches(1.8), Inches(5.6), Inches(0.5),
         "🎲 동작 방식", size=14, bold=True, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)

r_steps = [
    "① pc_url_pool.json 읽기 (크롤링이 수집한 풀)",
    "② 누적 이력(tested_urls.json) 자동 제외",
    "③ Fisher-Yates 셔플 → 200개 랜덤 샘플링",
    "④ 각 URL에 HTTP 200 + Body 404 체크",
    "⑤ 테스트 후 누적 이력에 추가",
    "⑥ 남은 URL < 100개 → 이력 초기화 (사이클 리셋)",
]
for i, text in enumerate(r_steps):
    add_text(s, Inches(0.8), Inches(2.5 + i * 0.6), Inches(5.5), Inches(0.5),
             text, size=13, color=TEXT_MAIN)

# 오른쪽: 다이어그램
add_rect(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.8), fill=BG_CARD, line=BORDER)
add_text(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(0.4),
         "📊 커버리지 확장 방식", size=14, bold=True, color=PURPLE)

# Day 1, 2, 3 박스
days = [
    ("Day 1", "Crawling 100 + Random 200", "누적 300 URL"),
    ("Day 2", "기존 300 제외 + 새 200", "누적 500 URL"),
    ("Day 3", "기존 500 제외 + 새 200", "누적 700 URL"),
    ("...", "남은 풀 < 100 → 사이클 리셋", "새 주기 시작"),
]
for i, (day, desc1, desc2) in enumerate(days):
    y = Inches(2.5 + i * 0.85)
    color = [PURPLE, ACCENT, PRIMARY, TEXT_HINT][i]
    add_rect(s, Inches(7.0), y, Inches(1.0), Inches(0.7), fill=color)
    add_text(s, Inches(7.0), y, Inches(1.0), Inches(0.7),
             day, size=13, bold=True, color=BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.2), y + Inches(0.05), Inches(4.5), Inches(0.3),
             desc1, size=11, bold=True, color=TEXT_MAIN)
    add_text(s, Inches(8.2), y + Inches(0.35), Inches(4.5), Inches(0.3),
             desc2, size=10, color=TEXT_SUB)

footer(s, next_page())


# ═══════════════════════════════════════════════════════════════
# 슬라이드 7: AutoE2E Test — 개요
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
header_bar(s, "06 · AutoE2E Test",
           "실제 사용자 시나리오 기반 기능 회귀 검증 (140+ TC)")

# 파일별 카드 테이블
headers = ["파일", "테스트", "우선순위", "주요 내용"]
rows_data = [
    ("auth.spec.ts",       "5",  "critical", "로그인·로그아웃·실패·회원가입 버튼"),
    ("cart.spec.ts",       "4",  "critical", "장바구니 진입·결제예정금액·카운트 증가"),
    ("product.spec.ts",    "8",  "critical", "상품명·가격·구매·옵션 레이어·쿠폰"),
    ("mypage.spec.ts",     "14", "high",     "주문/배송·상담·문의·회원정보·배송지"),
    ("search.spec.ts",     "10", "high",     "키워드·정렬/필터·영문·0건·탭"),
    ("navigation.spec.ts", "9",  "high",     "GNB 9개 카테고리 + 메인 홈"),
    ("category.spec.ts",   "13", "medium",   "9개 카테고리 + 필터 칩 + 섹션"),
    ("furnishing.spec.ts", "6",  "medium",   "특가·베스트·신상·1분 홈투어"),
    ("interior.spec.ts",   "12", "medium",   "공간 필터·시공사례·후기·상담"),
    ("store.spec.ts",      "2",  "low",      "매장찾기 페이지 진입"),
]

col_w = [Inches(2.4), Inches(0.9), Inches(1.5), Inches(7.2)]
table_x = Inches(0.6)
# header
x = table_x
for i, (hdr, w) in enumerate(zip(headers, col_w)):
    add_rect(s, x, Inches(1.8), w, Inches(0.45), fill=ACCENT)
    add_text(s, x + Inches(0.12), Inches(1.8), w - Inches(0.24), Inches(0.45),
             hdr, size=11, bold=True, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)
    x += w

priority_colors = {"critical": RED, "high": AMBER, "medium": PRIMARY, "low": TEXT_SUB}
# rows
for r_i, row in enumerate(rows_data):
    x = table_x
    y = Inches(2.25 + r_i * 0.42)
    bg = BG_CARD if r_i % 2 == 0 else BG_LIGHT
    for c_i, (val, w) in enumerate(zip(row, col_w)):
        add_rect(s, x, y, w, Inches(0.42), fill=bg, line=BORDER)
        if c_i == 2:
            # priority badge
            bc = priority_colors[val]
            add_rect(s, x + Inches(0.25), y + Inches(0.08), w - Inches(0.5), Inches(0.27), fill=bc)
            add_text(s, x + Inches(0.25), y + Inches(0.08), w - Inches(0.5), Inches(0.27),
                     val, size=9, bold=True, color=BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            color_v = TEXT_MAIN if c_i == 0 else TEXT_SUB
            bold_v = (c_i == 0)
            add_text(s, x + Inches(0.12), y, w - Inches(0.24), Inches(0.42),
                     val, size=10, color=color_v, bold=bold_v, anchor=MSO_ANCHOR.MIDDLE)
        x += w

# 합계 박스
add_rect(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.55), fill=PRIMARY)
add_text(s, Inches(0.8), Inches(6.6), Inches(11.9), Inches(0.55),
         "📊  합계:  PC E2E 약 83건 (Public_PC + Authed)   ·   MW E2E 약 52건 (isMobile skip 20~30)",
         size=13, bold=True, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)

footer(s, next_page())


# ═══════════════════════════════════════════════════════════════
# 슬라이드 8: POM & 제외 영역
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
header_bar(s, "07 · POM · 자동화 제외 영역", "Page Object Model + 파괴적 시나리오 식별")

# 왼쪽: POM
add_rect(s, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), fill=BG_CARD, line=BORDER)
add_rect(s, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.5), fill=TEAL)
add_text(s, Inches(0.7), Inches(1.8), Inches(5.6), Inches(0.5),
         "🏗 Page Object Model", size=14, bold=True, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.7), Inches(2.5), Inches(5.6), Inches(0.35),
         "셀렉터 우선순위:", size=12, bold=True, color=TEXT_MAIN)
add_text(s, Inches(0.7), Inches(2.9), Inches(5.6), Inches(0.35),
         "data-testid > aria-label > role+name > id > text > css",
         size=11, color=PRIMARY, bold=True)

poms = [
    ("BasePage",       "STORE_BASE · MALL_BASE · isMobile · smartLocator"),
    ("NavigationPage", "GNB 링크 · 검색창 · 로그인/장바구니 엔트리"),
    ("CategoryPage",   "CATEGORIES 상수 · 필터 칩 · firstGoodsLink"),
    ("CartPage",       "장바구니 URL · getCount() · isEmpty()"),
]
for i, (name, desc) in enumerate(poms):
    y = Inches(3.5 + i * 0.75)
    add_rect(s, Inches(0.7), y, Inches(5.6), Inches(0.65), fill=BG_LIGHT)
    add_text(s, Inches(0.9), y + Inches(0.05), Inches(5.2), Inches(0.3),
             name, size=12, bold=True, color=TEAL)
    add_text(s, Inches(0.9), y + Inches(0.32), Inches(5.2), Inches(0.3),
             desc, size=10, color=TEXT_SUB)

# 오른쪽: 제외 영역
add_rect(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0), fill=BG_CARD, line=BORDER)
add_rect(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(0.5), fill=RED)
add_text(s, Inches(7.0), Inches(1.8), Inches(5.6), Inches(0.5),
         "🚫 자동화 불가 · 제외 영역", size=14, bold=True, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)

exclusions = [
    ("💳", "실결제 7종",          "카드·네이버·카카오·토스·페이코·퀵·휴대폰"),
    ("👤", "회원가입 / 탈퇴",      "파괴적 작업"),
    ("🔐", "OAuth 간편가입",      "외부 IdP (카카오·네이버·애플)"),
    ("📦", "주문 취소 / 반품",    "실주문 전제"),
    ("📧", "ID / PW 찾기",        "이메일·SMS 인증 필요"),
]
for i, (icon, title, desc) in enumerate(exclusions):
    y = Inches(2.5 + i * 0.85)
    add_text(s, Inches(7.0), y + Inches(0.1), Inches(0.5), Inches(0.5),
             icon, size=20)
    add_text(s, Inches(7.6), y + Inches(0.05), Inches(5.1), Inches(0.35),
             title, size=13, bold=True, color=TEXT_MAIN)
    add_text(s, Inches(7.6), y + Inches(0.42), Inches(5.1), Inches(0.35),
             desc, size=10, color=TEXT_SUB)

footer(s, next_page())


# ═══════════════════════════════════════════════════════════════
# 슬라이드 9: 보조 시스템 (AI 분석·커버리지·알림)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
header_bar(s, "08 · 보조 시스템", "AI 실패 분석 · 커버리지 · 통합 알림")

# 3개 카드
cards = [
    {
        "icon": "🤖",
        "title": "Claude AI 실패 분석",
        "color": TEAL,
        "items": [
            "실패 케이스 Claude Sonnet 4.6 전송",
            "JSON 응답 파싱:",
            "  summary (한 줄 요약)",
            "  rootCauses (원인 목록)",
            "  suggestedFixes (수정 제안)",
            "  priority (critical/high/medium/low)",
            "detail.html 'AI 분석' 박스에 표시",
            "ANTHROPIC_API_KEY 없으면 skip",
        ],
    },
    {
        "icon": "📈",
        "title": "커버리지 통계",
        "color": PURPLE,
        "items": [
            "누적 고유 URL 수 (PC+MW)",
            "오늘 테스트한 URL 수",
            "오늘 통과율",
            "완주 사이클 수",
            "실시간 대시보드 상단 바 표시",
            "pc_tested_urls.json 기반",
            "매 테스트 후 자동 갱신",
            "coverage.ts에서 집계 로직",
        ],
    },
    {
        "icon": "💬",
        "title": "통합 잔디 알림",
        "color": ACCENT,
        "items": [
            "모든 테스트 종료 후 1회만 전송",
            "6개 리포트별 통과율 요약",
            "전체 총계 + 통과율",
            "실패 URL 상위 10건",
            "대시보드 링크",
            "실패 0 → 초록 · 있으면 빨강",
            "scripts/send-jandi-summary.js",
            "Actions if: always() → 실패 시에도 전송",
        ],
    },
]
card_w = Inches(4.0)
card_h = Inches(5.0)
for i, card in enumerate(cards):
    x = Inches(0.5 + i * 4.2)
    add_rect(s, x, Inches(1.8), card_w, card_h, fill=BG_CARD, line=BORDER)
    add_rect(s, x, Inches(1.8), card_w, Inches(0.6), fill=card["color"])
    add_text(s, x + Inches(0.2), Inches(1.85), Inches(0.5), Inches(0.5),
             card["icon"], size=20, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.9), Inches(1.85), card_w - Inches(1.0), Inches(0.5),
             card["title"], size=14, bold=True, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)
    for j, text in enumerate(card["items"]):
        add_text(s, x + Inches(0.3), Inches(2.6 + j * 0.35), card_w - Inches(0.6), Inches(0.3),
                 ("▸ " if not text.startswith("  ") else "   ") + text.strip(),
                 size=10.5, color=TEXT_SUB)

footer(s, next_page())


# ═══════════════════════════════════════════════════════════════
# 슬라이드 10: 대시보드 UI
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
header_bar(s, "09 · 대시보드 UI", "index.html · detail.html · incidents.html")

# 대시보드 모형: 커버리지 바 + 6개 카드
add_rect(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.8), fill=BG_CARD, line=BORDER)
add_text(s, Inches(0.7), Inches(2.15), Inches(12.0), Inches(0.25),
         "📊 커버리지 통계 바",
         size=11, bold=True, color=TEXT_HINT)
stats = [
    ("누적 URL", "1,819", TEXT_MAIN),
    ("PC / MW", "914 / 905", TEXT_SUB),
    ("오늘 테스트", "600건", TEXT_SUB),
    ("통과율", "99.5%", GREEN),
    ("사이클", "0 / 0", TEXT_SUB),
]
sx = Inches(0.9)
for i, (label, value, color) in enumerate(stats):
    add_text(s, sx + Inches(i * 2.4), Inches(2.45), Inches(2.3), Inches(0.3),
             label, size=9, color=TEXT_HINT)
    add_text(s, sx + Inches(i * 2.4), Inches(2.55), Inches(2.3), Inches(0.3),
             value, size=15, bold=True, color=color)

# 6개 카드
card_configs = [
    ("PC 500 Crawling", "CRAWL", PRIMARY, "100%"),
    ("PC 200 Random", "RANDOM", PURPLE, "100%"),
    ("PC E2E Test", "E2E", ACCENT, "93.8%"),
    ("MW 500 Crawling", "CRAWL", TEAL, "100%"),
    ("MW 200 Random", "RANDOM", RGBColor(0xAD, 0x14, 0x57), "100%"),
    ("MW E2E Test", "E2E", GREEN, "95.0%"),
]
cw = Inches(3.9); ch = Inches(1.7); gap = Inches(0.15)
for i, (title, label, color, rate) in enumerate(card_configs):
    r, c = divmod(i, 3)
    x = Inches(0.5) + c * (cw + gap)
    y = Inches(3.1) + r * (ch + Inches(0.15))
    add_rect(s, x, y, cw, ch, fill=BG_CARD, line=BORDER)
    # icon box
    add_rect(s, x + Inches(0.15), y + Inches(0.2), Inches(0.7), Inches(0.7), fill=color)
    # abbr1
    add_text(s, x + Inches(0.15), y + Inches(0.22), Inches(0.7), Inches(0.3),
             title.split()[0], size=10, bold=True, color=BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # abbr2
    add_text(s, x + Inches(0.15), y + Inches(0.55), Inches(0.7), Inches(0.3),
             label, size=9, color=BG_CARD, align=PP_ALIGN.CENTER)
    # title
    add_text(s, x + Inches(1.0), y + Inches(0.2), cw - Inches(1.2), Inches(0.4),
             title, size=12, bold=True, color=TEXT_MAIN, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.0), y + Inches(0.55), cw - Inches(1.2), Inches(0.3),
             "최근 실행: 오늘", size=9, color=TEXT_HINT)
    # rate
    rate_color = GREEN if "100" in rate or float(rate.replace("%", "")) >= 95 else AMBER
    add_text(s, x + Inches(0.2), y + Inches(1.0), cw - Inches(0.4), Inches(0.6),
             rate, size=28, bold=True, color=rate_color, align=PP_ALIGN.RIGHT)

footer(s, next_page())


# ═══════════════════════════════════════════════════════════════
# 슬라이드 11: 향후 개선 — 단기 (1~2주)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
header_bar(s, "10 · 향후 개선 (단기) · 1~2주", "즉시 착수 가능 · 운영 안정화")

tasks = [
    ("🔴", "HIGH", "MAX_LINKS 원복 (100 → 500)", "실행 시간 분석 후 단계적 확장", RED),
    ("🔴", "HIGH", "실패 케이스 TC 보강", "검색 결과·필터·상품 상세 선택기 안정화", RED),
    ("🔴", "HIGH", "잔디 알림 검증", "다음 07:00 자동 실행 후 통합 알림 포맷 확인", RED),
    ("🟡", "MED",  "CLAUDE.md 업데이트", "autoe2e 폴더 · AI 분석 · 커버리지 반영", AMBER),
    ("🟡", "MED",  "로컬 개발 가이드", "npm test · --grep · 디버깅 팁 문서화", AMBER),
]

for i, (icon, pri, title, desc, color) in enumerate(tasks):
    y = Inches(1.9 + i * 0.95)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.85), fill=BG_CARD, line=BORDER)
    add_rect(s, Inches(0.5), y, Inches(0.12), Inches(0.85), fill=color)
    add_text(s, Inches(0.8), y, Inches(0.6), Inches(0.85),
             icon, size=24, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(1.5), y + Inches(0.27), Inches(0.8), Inches(0.3), fill=color)
    add_text(s, Inches(1.5), y + Inches(0.27), Inches(0.8), Inches(0.3),
             pri, size=10, bold=True, color=BG_CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.5), y + Inches(0.05), Inches(10.2), Inches(0.35),
             title, size=14, bold=True, color=TEXT_MAIN)
    add_text(s, Inches(2.5), y + Inches(0.45), Inches(10.2), Inches(0.35),
             desc, size=11, color=TEXT_SUB)

footer(s, next_page())


# ═══════════════════════════════════════════════════════════════
# 슬라이드 12: 향후 개선 — 중기 (1~2개월)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
header_bar(s, "11 · 향후 개선 (중기) · 1~2개월", "관찰성 강화 + 테스트 확장")

# 2 컬럼 레이아웃
# 좌: 관찰성
add_rect(s, Inches(0.5), Inches(1.9), Inches(6.0), Inches(5.0), fill=BG_CARD, line=BORDER)
add_rect(s, Inches(0.5), Inches(1.9), Inches(6.0), Inches(0.6), fill=PURPLE)
add_text(s, Inches(0.7), Inches(1.9), Inches(5.6), Inches(0.6),
         "📊 관찰성 강화", size=15, bold=True, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)

obs_items = [
    ("주간·월간 추이 그래프", "results.json 히스토리 시계열 분석"),
    ("Slack/Teams 알림 확장", "팀별 선호 채널 추가"),
    ("실패 URL 반복 감지", "동일 URL 연속 실패 시 에스컬레이션"),
    ("Grafana 연동 준비", "JSON → DB 이관 검토"),
]
for i, (t, d) in enumerate(obs_items):
    y = Inches(2.65 + i * 1.05)
    add_text(s, Inches(0.8), y, Inches(5.5), Inches(0.4),
             f"✓ {t}", size=13, bold=True, color=TEXT_MAIN)
    add_text(s, Inches(1.1), y + Inches(0.42), Inches(5.2), Inches(0.4),
             d, size=11, color=TEXT_SUB)

# 우: 테스트 확장
add_rect(s, Inches(6.8), Inches(1.9), Inches(6.0), Inches(5.0), fill=BG_CARD, line=BORDER)
add_rect(s, Inches(6.8), Inches(1.9), Inches(6.0), Inches(0.6), fill=GREEN)
add_text(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(0.6),
         "🧪 테스트 확장", size=15, bold=True, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)

ext_items = [
    ("Visual Regression", "Playwright 내장 스크린샷 비교"),
    ("Core Web Vitals 측정", "LCP · FID · CLS 자동 수집"),
    ("접근성 테스트", "@axe-core/playwright 통합"),
    ("API 레벨 스모크", "페이지 로딩 없이 엔드포인트 직접 호출"),
]
for i, (t, d) in enumerate(ext_items):
    y = Inches(2.65 + i * 1.05)
    add_text(s, Inches(7.1), y, Inches(5.5), Inches(0.4),
             f"✓ {t}", size=13, bold=True, color=TEXT_MAIN)
    add_text(s, Inches(7.4), y + Inches(0.42), Inches(5.2), Inches(0.4),
             d, size=11, color=TEXT_SUB)

footer(s, next_page())


# ═══════════════════════════════════════════════════════════════
# 슬라이드 13: 향후 개선 — 장기 (3~6개월)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=BG_LIGHT)
header_bar(s, "12 · 향후 개선 (장기) · 3~6개월", "AI 기반 자동화 + 인프라 고도화")

long_cards = [
    {
        "icon": "🤖", "title": "AI 기반 자동화 고도화", "color": TEAL,
        "items": [
            ("셀렉터 자동 수리",
             "실패 시 Claude에게 스크린샷+HTML 전달 → 새 셀렉터 제안 → PR 자동 생성"),
            ("테스트 자동 생성",
             "AutoE2E 파이프라인 재활용 → 새 페이지 감지 시 시나리오 자동 작성"),
            ("장애 예측 모델",
             "과거 실패 패턴 학습 → 위험도 높은 URL 사전 경고"),
        ],
    },
    {
        "icon": "🏗", "title": "인프라 고도화", "color": ACCENT,
        "items": [
            ("Self-hosted Runner",
             "내부망 접근 가능 → Stage 환경 테스트 재개"),
            ("Parallel Execution",
             "workers: 1 제약 완화 + 요청 간 지연 → 실행 시간 50% 단축"),
            ("결과 DB 저장",
             "JSON → PostgreSQL/SQLite → Grafana 대시보드 연동"),
        ],
    },
]
for ci, card in enumerate(long_cards):
    x = Inches(0.5 + ci * 6.4)
    add_rect(s, x, Inches(1.9), Inches(6.2), Inches(5.0), fill=BG_CARD, line=BORDER)
    add_rect(s, x, Inches(1.9), Inches(6.2), Inches(0.6), fill=card["color"])
    add_text(s, x + Inches(0.2), Inches(1.9), Inches(0.4), Inches(0.6),
             card["icon"], size=20, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.8), Inches(1.9), Inches(5.3), Inches(0.6),
             card["title"], size=15, bold=True, color=BG_CARD, anchor=MSO_ANCHOR.MIDDLE)
    for i, (t, d) in enumerate(card["items"]):
        y = Inches(2.7 + i * 1.35)
        add_rect(s, x + Inches(0.15), y, Inches(5.9), Inches(1.2), fill=BG_LIGHT)
        add_text(s, x + Inches(0.35), y + Inches(0.1), Inches(5.5), Inches(0.35),
                 f"▸ {t}", size=13, bold=True, color=card["color"])
        add_text(s, x + Inches(0.6), y + Inches(0.5), Inches(5.3), Inches(0.65),
                 d, size=11, color=TEXT_SUB)

footer(s, next_page())


# ═══════════════════════════════════════════════════════════════
# 슬라이드 14: 마무리
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=PRIMARY)

add_text(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.6),
         "THANK YOU", size=48, bold=True, color=BG_CARD, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.5),
         "한샘몰 서비스 품질 관리 대시보드", size=22, color=BG_CARD, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.4),
         "자동 테스트로 사용자 경험을 지킵니다.", size=14, color=RGBColor(0xE0, 0xE5, 0xF0), align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.3),
         "📍 daewon82.github.io/hanssem-qa-system",
         size=12, color=RGBColor(0xE0, 0xE5, 0xF0), align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.3),
         "📂 github.com/daewon82/hanssem-qa-system",
         size=12, color=RGBColor(0xE0, 0xE5, 0xF0), align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.3),
         "dwlee@hanssem.com · 2026-04-23",
         size=10, color=RGBColor(0xC0, 0xC8, 0xDB), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════════════
output = "/Users/dw/hanssem-qa-system/한샘몰_QA시스템_소개.pptx"
prs.save(output)
print(f"✅ PPT 생성 완료: {output}")
print(f"   총 {len(prs.slides)}개 슬라이드")
