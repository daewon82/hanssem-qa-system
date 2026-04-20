from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── 색상 팔레트 ────────────────────────────────────────────────
HANSSEM_BLUE   = RGBColor(0x00, 0x5B, 0xAC)   # 한샘 블루
HANSSEM_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # 다크 네이비
ACCENT_ORANGE  = RGBColor(0xFF, 0x6B, 0x35)   # 강조 오렌지
LIGHT_BLUE     = RGBColor(0xE8, 0xF4, 0xFF)   # 연한 파랑
LIGHT_GRAY     = RGBColor(0xF5, 0xF5, 0xF5)   # 연한 회색
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT      = RGBColor(0x2C, 0x2C, 0x2C)
MID_GRAY       = RGBColor(0x6B, 0x6B, 0x6B)
GREEN          = RGBColor(0x00, 0xB3, 0x5A)
RED            = RGBColor(0xE5, 0x39, 0x35)
YELLOW         = RGBColor(0xFF, 0xC1, 0x07)
TABLE_HEADER   = RGBColor(0x00, 0x5B, 0xAC)
TABLE_ROW_ODD  = RGBColor(0xF0, 0xF6, 0xFF)
TABLE_ROW_EVEN = RGBColor(0xFF, 0xFF, 0xFF)
COVER_GRAD     = RGBColor(0x00, 0x3A, 0x7D)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 완전 빈 레이아웃

# ── 헬퍼 함수 ─────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = Pt(line_w)
    else:
        if not line:
            shape.line.fill.background()
    return shape

def add_text_box(slide, text, x, y, w, h,
                 font_size=14, bold=False, color=DARK_TEXT,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"
    return txBox

def add_multiline(slide, lines, x, y, w, h,
                  font_size=12, bold=False, color=DARK_TEXT,
                  align=PP_ALIGN.LEFT, line_spacing=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as PT
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "맑은 고딕"
    return txBox

def slide_header(slide, title, subtitle=None, bar_color=HANSSEM_BLUE):
    # 상단 컬러 바
    add_rect(slide, 0, 0, 13.33, 1.1, fill=bar_color)
    # 슬라이드 번호 장식 원
    add_rect(slide, 0.25, 0.15, 0.75, 0.75, fill=ACCENT_ORANGE)
    add_text_box(slide, title, 1.2, 0.12, 10.5, 0.85,
                 font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 1.2, 0.8, 10.5, 0.4,
                     font_size=11, color=RGBColor(0xCC, 0xE4, 0xFF), align=PP_ALIGN.LEFT)

def add_table(slide, headers, rows, x, y, w, h, col_widths=None):
    cols = len(headers)
    nrows = len(rows) + 1
    table = slide.shapes.add_table(nrows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table

    if col_widths:
        total_w = sum(col_widths)
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = Inches(w * cw / total_w)

    def set_cell(cell, text, bg, font_color=WHITE, bold=False, sz=10.5, align=PP_ALIGN.CENTER):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = font_color
        run.font.name = "맑은 고딕"
        from pptx.util import Pt as PT
        cell.margin_top = PT(4)
        cell.margin_bottom = PT(4)
        cell.margin_left = PT(6)
        cell.margin_right = PT(6)

    for ci, h_txt in enumerate(headers):
        set_cell(table.cell(0, ci), h_txt, TABLE_HEADER, WHITE, bold=True, sz=10.5)

    for ri, row in enumerate(rows):
        bg = TABLE_ROW_ODD if ri % 2 == 0 else TABLE_ROW_EVEN
        for ci, val in enumerate(row):
            align = PP_ALIGN.LEFT if ci == 0 or (isinstance(val, str) and len(val) > 10) else PP_ALIGN.CENTER
            set_cell(table.cell(ri+1, ci), val, bg, DARK_TEXT, sz=10, align=align)

    return table

def badge(slide, text, x, y, w=1.2, h=0.32, bg=HANSSEM_BLUE, fg=WHITE, sz=10):
    add_rect(slide, x, y, w, h, fill=bg)
    add_text_box(slide, text, x, y+0.01, w, h-0.02, font_size=sz, bold=True,
                 color=fg, align=PP_ALIGN.CENTER)

def section_box(slide, title, x, y, w, h, title_color=HANSSEM_BLUE):
    add_rect(slide, x, y, w, h, fill=LIGHT_BLUE, line=HANSSEM_BLUE, line_w=0.5)
    add_rect(slide, x, y, w, 0.35, fill=title_color)
    add_text_box(slide, title, x+0.1, y+0.02, w-0.2, 0.31,
                 font_size=11, bold=True, color=WHITE)

# ══════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# 배경
add_rect(slide, 0, 0, 13.33, 7.5, fill=HANSSEM_DARK)
# 하단 강조 줄
add_rect(slide, 0, 6.8, 13.33, 0.7, fill=HANSSEM_BLUE)
# 오른쪽 장식
add_rect(slide, 10.5, 0, 2.83, 7.5, fill=RGBColor(0x00, 0x3A, 0x7D))
add_rect(slide, 10.5, 0, 0.06, 7.5, fill=ACCENT_ORANGE)

# 왼쪽 콘텐츠
add_text_box(slide, "HANSSEM MALL", 0.6, 1.0, 9.5, 0.6,
             font_size=14, bold=False, color=RGBColor(0x88, 0xBB, 0xFF),
             align=PP_ALIGN.LEFT)
add_text_box(slide, "서비스 품질 관리 시스템", 0.6, 1.55, 9.5, 1.0,
             font_size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(slide, "QA System Technical Document", 0.6, 2.5, 9.5, 0.5,
             font_size=16, bold=False, color=RGBColor(0xAA, 0xCC, 0xFF),
             align=PP_ALIGN.LEFT)

# 구분선
add_rect(slide, 0.6, 3.1, 4.0, 0.04, fill=ACCENT_ORANGE)

# 설명
add_multiline(slide, [
    "Playwright + TypeScript + GitHub Actions + GitHub Pages",
    "PC / MW 운영환경 자동 QA — 랜딩 500개 + E2E 시나리오 테스트",
], 0.6, 3.25, 9.5, 0.9, font_size=13, color=RGBColor(0xCC, 0xDD, 0xFF))

# 메타
add_text_box(slide, "작성일: 2026-04-20", 0.6, 6.85, 5.0, 0.4,
             font_size=11, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(slide, "대상: 개발팀 / QA 담당자", 5.5, 6.85, 5.0, 0.4,
             font_size=11, color=WHITE, align=PP_ALIGN.LEFT)

# 오른쪽 아이콘 영역
icons = ["🔍 랜딩 테스트", "🖱️ E2E 테스트", "📊 대시보드", "🔔 잔디 알림"]
for i, ic in enumerate(icons):
    add_text_box(slide, ic, 10.65, 1.5 + i*1.1, 2.5, 0.6,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 슬라이드 2: 목차
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(slide, "목차", "Table of Contents")

toc = [
    ("01", "시스템 개요 및 기대효과", "도입 배경 · 효과 · 구성 요약"),
    ("02", "전체 아키텍처", "실행 흐름 · 디렉토리 구조 · Playwright 설정"),
    ("03", "공통 헬퍼 (utils.ts)", "updateProgress · publishResults"),
    ("04", "테스트 구성 상세", "랜딩 500개 · PC E2E 18그룹 · MW E2E 13그룹"),
    ("05", "결과 데이터 & 대시보드", "results.json 구조 · 실시간 로딩 UI"),
    ("06", "GitHub Actions 워크플로우", "daily-trigger · playwright.yml 실행 단계"),
    ("07", "기능 커버리지 분석", "66개 기본기능 대비 약 21% · 미커버 원인 분석"),
    ("08", "운영 가이드", "장애 리포트 추가 · 잔디 알림 · PAT 토큰 관리"),
]

for i, (num, title, sub) in enumerate(toc):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.3 + row * 1.45

    add_rect(slide, x, y, 6.1, 1.2, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
    add_rect(slide, x, y, 0.7, 1.2, fill=HANSSEM_BLUE)
    add_text_box(slide, num, x, y+0.3, 0.7, 0.5,
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, x+0.8, y+0.08, 5.1, 0.5,
                 font_size=13, bold=True, color=HANSSEM_DARK, align=PP_ALIGN.LEFT)
    add_text_box(slide, sub, x+0.8, y+0.58, 5.1, 0.5,
                 font_size=10, color=MID_GRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════
# 슬라이드 3: 시스템 개요 및 기대효과
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(slide, "01. 시스템 개요 및 기대효과", "System Overview & Expected Effects")

# 도입 배경
add_rect(slide, 0.3, 1.25, 12.73, 1.45, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 0.3, 1.25, 0.08, 1.45, fill=ACCENT_ORANGE)
add_text_box(slide, "도입 배경", 0.55, 1.28, 3.0, 0.35, font_size=11, bold=True, color=HANSSEM_BLUE)
add_multiline(slide, [
    "한샘몰은 PC(store.hanssem.com)와 MW(m.store.hanssem.com) 두 채널을 운영하며,",
    "배포 후 수백 개의 URL과 주요 사용자 시나리오를 수동으로 확인하는 데 많은 리소스가 필요했다.",
    "이를 자동화해 배포 익일 아침 이상을 조기 감지하고, QA 리소스를 절감하는 것이 목표이다.",
], 0.55, 1.6, 12.3, 1.0, font_size=11, color=DARK_TEXT)

# 기대효과 표
headers = ["항목", "도입 전", "도입 후"]
rows = [
    ("랜딩 점검",     "수동 접속 확인 (PC+MW 각 500개)",   "매일 KST 08:00 자동 실행"),
    ("결과 확인",     "담당자 직접 기록·공유",              "대시보드에서 즉시 확인"),
    ("이상 감지",     "사용자 신고 후 인지",                "배포 익일 아침 자동 감지"),
    ("장애 이력",     "별도 문서 관리",                     "대시보드 통합 관리"),
    ("알림",          "수동 공유",                          "잔디 웹훅 자동 전송"),
]
add_table(slide, headers, rows, 0.3, 2.85, 8.0, 2.5, col_widths=[2.5, 3.5, 3.5])

# 구성 요약 박스
add_rect(slide, 8.6, 2.85, 4.4, 2.5, fill=HANSSEM_BLUE)
add_text_box(slide, "시스템 구성", 8.75, 2.9, 4.1, 0.35, font_size=11, bold=True, color=WHITE)
items = [
    "🔍 랜딩 테스트",
    "   PC 500개 + MW 500개 URL",
    "   HTTP 200 응답 자동 점검",
    "",
    "🖱️  E2E 테스트",
    "   PC 18개 + MW 13개 시나리오",
    "   사용자 동선 자동 재현",
    "",
    "📊 GitHub Pages 대시보드",
    "   결과 실시간 시각화",
]
add_multiline(slide, items, 8.75, 3.3, 4.1, 2.0, font_size=10, color=WHITE)

# 하단 키 수치
kpis = [("500개 × 2채널", "랜딩 URL 점검"), ("31개 시나리오", "E2E TC"), ("매일 08:00", "자동 실행"), ("약 21%", "기능 커버리지")]
for i, (val, lbl) in enumerate(kpis):
    x = 0.3 + i * 3.2
    add_rect(slide, x, 5.55, 3.0, 1.7, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
    add_text_box(slide, val, x, 5.75, 3.0, 0.7,
                 font_size=20, bold=True, color=HANSSEM_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, lbl, x, 6.45, 3.0, 0.4,
                 font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 슬라이드 4: 전체 아키텍처
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(slide, "02. 전체 아키텍처", "System Architecture")

# 흐름도 배경
add_rect(slide, 0.3, 1.2, 12.73, 3.7, fill=WHITE, line=RGBColor(0xDD,0xDD,0xDD), line_w=0.3)

# 트리거
add_rect(slide, 0.5, 1.4, 2.2, 0.9, fill=ACCENT_ORANGE)
add_multiline(slide, ["daily-trigger.yml", "KST 08:00 자동 실행"], 0.5, 1.4, 2.2, 0.9,
              font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "→", 2.75, 1.65, 0.4, 0.4, font_size=18, bold=True, color=HANSSEM_BLUE, align=PP_ALIGN.CENTER)

# playwright.yml
add_rect(slide, 3.2, 1.4, 2.5, 0.9, fill=HANSSEM_BLUE)
add_multiline(slide, ["playwright.yml", "workflow_dispatch"], 3.2, 1.4, 2.5, 0.9,
              font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 화살표 아래로
add_text_box(slide, "↓", 4.25, 2.35, 0.5, 0.4, font_size=18, bold=True, color=HANSSEM_BLUE, align=PP_ALIGN.CENTER)

# PC / MW 박스
for ci, (label, color, tests) in enumerate([
    ("PC_Chrome", HANSSEM_BLUE, ["pc_500_check.spec.ts  →  랜딩 500개", "pc_e2e_test_v2.spec.ts  →  E2E 18그룹"]),
    ("MW_Chrome", RGBColor(0x00, 0x80, 0xC0), ["mw_500_check.spec.ts  →  랜딩 500개", "mw_e2e_test.spec.ts     →  E2E 13그룹"]),
]):
    x = 0.5 + ci * 6.3
    add_rect(slide, x, 2.85, 5.9, 0.4, fill=color)
    add_text_box(slide, label, x, 2.88, 5.9, 0.34,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, t in enumerate(tests):
        add_rect(slide, x, 3.3 + j*0.55, 5.9, 0.45, fill=TABLE_ROW_ODD if j%2==0 else WHITE, line=color, line_w=0.3)
        add_text_box(slide, t, x+0.15, 3.33 + j*0.55, 5.6, 0.38, font_size=10, color=DARK_TEXT)

# 결과 흐름
add_text_box(slide, "↓ 각 테스트 완료 시 publishResults() → gh-pages 즉시 반영", 0.5, 4.5, 12.0, 0.35,
             font_size=10, color=HANSSEM_BLUE, bold=True)

# 하단 3단계
steps = [
    ("4단계", "결과 커밋", "main 브랜치\ngit pull --rebase && push"),
    ("5단계", "Pages 배포", "public/ → gh-pages\npeaciris/actions-gh-pages@v3"),
    ("6단계", "아티팩트", "playwright-report\nfail_evidence (30일 보관)"),
]
for i, (n, t, d) in enumerate(steps):
    x = 0.3 + i * 4.3
    add_rect(slide, x, 4.95, 4.0, 1.3, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
    add_rect(slide, x, 4.95, 4.0, 0.38, fill=HANSSEM_BLUE)
    add_text_box(slide, f"{n}  {t}", x+0.1, 4.98, 3.8, 0.32, font_size=10, bold=True, color=WHITE)
    add_text_box(slide, d, x+0.15, 5.38, 3.7, 0.82, font_size=10, color=DARK_TEXT)

# 디렉토리 구조 (우측 하단)
add_rect(slide, 12.7-0.01, 1.2, 0.08, 3.7, fill=ACCENT_ORANGE)

# ══════════════════════════════════════════════════════════════
# 슬라이드 5: 공통 헬퍼 utils.ts
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(slide, "03. 공통 헬퍼 — utils.ts", "Shared Helpers")

# updateProgress
add_rect(slide, 0.3, 1.25, 6.1, 5.9, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 0.3, 1.25, 6.1, 0.42, fill=HANSSEM_BLUE)
add_text_box(slide, "updateProgress(phase, count?, total?)", 0.45, 1.28, 5.8, 0.36,
             font_size=11, bold=True, color=WHITE)

up_lines = [
    "실행 단계를 gh-pages의 progress.json에 실시간 기록",
    "대시보드가 15초마다 폴링 → 로딩 UI 및 진행률 표시",
    "",
    "파라미터",
    "  phase : 'pc-landing' | 'pc-e2e' | 'mw-landing' | 'mw-e2e'",
    "  count : 현재 방문 건수 (랜딩 테스트)",
    "  total : 전체 건수 (랜딩 테스트)",
    "",
    "동작 조건",
    "  GITHUB_ACTIONS + GITHUB_TOKEN 환경변수 필요",
    "  로컬 실행 시 자동 스킵 (안전)",
    "  gh-pages 브랜치에 직접 PUT (main 커밋 없음)",
    "  실패해도 경고 로그만 출력 → 테스트 계속 진행",
    "",
    "progress.json 구조",
    '  { "phase": "pc-landing",',
    '    "startedAt": "2026-04-20T00:00:00.000Z",',
    '    "count": 150, "total": 500 }',
]
add_multiline(slide, up_lines, 0.45, 1.75, 5.8, 5.0, font_size=10, color=DARK_TEXT)

# publishResults
add_rect(slide, 6.7, 1.25, 6.3, 5.9, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 6.7, 1.25, 6.3, 0.42, fill=RGBColor(0x00, 0x80, 0xC0))
add_text_box(slide, "publishResults(report, fullData, fullDataPath)", 6.85, 1.28, 6.0, 0.36,
             font_size=11, bold=True, color=WHITE)

pr_lines = [
    "각 테스트 afterAll에서 호출",
    "결과를 gh-pages에 즉시 반영 → 실시간 대시보드 갱신",
    "",
    "파라미터",
    "  report     : 리포트 요약 객체",
    "               (id, title, total, pass, fail, passRate, cases)",
    "  fullData   : 전체 케이스 데이터 (pass + fail)",
    "  fullDataPath : 'pc_500.json' | 'pc_e2e.json' | ...",
    "",
    "동작 순서",
    "  1. fullDataPath 파일 gh-pages에 직접 PUT",
    "  2. results.json GET → report 병합 → PUT",
    "  3. 완료 로그 출력 (📊 결과 즉시 반영 완료)",
    "",
    "동작 조건",
    "  GITHUB_ACTIONS + GITHUB_TOKEN 환경변수 필요",
    "  로컬 실행 시 자동 스킵",
]
add_multiline(slide, pr_lines, 6.85, 1.75, 6.0, 5.0, font_size=10, color=DARK_TEXT)

# 구분선
add_rect(slide, 6.55, 1.25, 0.08, 5.9, fill=ACCENT_ORANGE)

# ══════════════════════════════════════════════════════════════
# 슬라이드 6: 랜딩 테스트
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(slide, "04-A. 랜딩 테스트 (PC / MW)", "Landing Test — 500 URLs per Channel")

# PC vs MW 비교표
headers = ["항목", "PC 랜딩", "MW 랜딩"]
rows = [
    ("파일",       "pc_500_check.spec.ts",         "mw_500_check.spec.ts"),
    ("시작 URL",   "store.hanssem.com",             "m.store.hanssem.com"),
    ("링크 필터",  "store.hanssem.com 포함",        "m.store.hanssem.com 포함"),
    ("점검 건수",  "500개",                          "500개"),
    ("결과 파일",  "public/pc_500.json",            "public/mw_500.json"),
    ("ID",        "pc-landing",                     "mw-landing"),
]
add_table(slide, headers, rows, 0.3, 1.25, 8.0, 2.5, col_widths=[2.5, 3.5, 3.5])

# 실패 판정 박스
add_rect(slide, 8.5, 1.25, 4.5, 2.5, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 8.5, 1.25, 4.5, 0.38, fill=HANSSEM_BLUE)
add_text_box(slide, "실패 판정 기준", 8.65, 1.28, 4.2, 0.32, font_size=11, bold=True, color=WHITE)
add_multiline(slide, [
    "HTTP 200  →  ✅ PASS",
    "HTTP 오류 / Timeout  →  ❌ FAIL",
    "",
    "5xx 응답: 5초 후 1회 재시도",
    "타임아웃 (연속 3회 미만): 20초 대기 후 재시도",
    "타임아웃 (연속 3회 이상): 60초 대기 후 재시도",
    "최종 실패: window.stop() → 스크린샷 저장",
], 8.65, 1.7, 4.2, 2.0, font_size=10, color=DARK_TEXT)

# 링크 수집 방식
add_rect(slide, 0.3, 3.9, 5.8, 3.35, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 0.3, 3.9, 5.8, 0.38, fill=RGBColor(0x00, 0x80, 0xC0))
add_text_box(slide, "링크 수집 방식", 0.45, 3.93, 5.5, 0.32, font_size=11, bold=True, color=WHITE)
add_multiline(slide, [
    "1. 시작 URL에서 <a href> 전체 수집",
    "   (스크롤 3회로 lazy-load 유발)",
    "2. URL 유효성 필터 적용",
    "   · 길이 400자 초과 제외",
    "   · 정적 파일 확장자 제외 (.pdf .jpg .css .js 등)",
    "   · API/정적 경로 제외 (/api/ /static/ /assets/)",
    "3. EXCLUDE_KEYWORDS 필터",
    "   logout, login, javascript, order, settle,",
    "   cart, member, company.hanssem.com",
    "4. 500개 채울 때까지 방문 페이지에서 추가 수집 반복",
    "5. 5초마다 updateProgress() 호출",
], 0.45, 4.35, 5.5, 2.8, font_size=10, color=DARK_TEXT)

# 구조 특이사항
add_rect(slide, 6.4, 3.9, 6.6, 3.35, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 6.4, 3.9, 6.6, 0.38, fill=ACCENT_ORANGE)
add_text_box(slide, "구조 특이사항 및 진행상태 업데이트", 6.55, 3.93, 6.3, 0.32, font_size=11, bold=True, color=WHITE)
add_multiline(slide, [
    "구조",
    "  · 테스트 함수 1개 안에서 while 루프 실행",
    "  · Playwright 카운트: 항상 '1 passed'",
    "  · 실제 pass/fail은 내부 변수로 집계 후 JSON 저장",
    "",
    "진행상태 업데이트",
    "  · 테스트 시작 시 1회: updateProgress('pc-landing', 0, 500)",
    "  · 이후 5초마다: updateProgress('pc-landing', count, 500)",
    "  · 대시보드가 진행률(%) 실시간 표시",
    "",
    "결과 저장",
    "  · results.json: fail 케이스만 저장",
    "  · pc_500.json: pass + fail 전체 저장",
    "  · afterAll 종료 시 publishResults() 호출",
], 6.55, 4.35, 6.3, 2.8, font_size=10, color=DARK_TEXT)

# ══════════════════════════════════════════════════════════════
# 슬라이드 7: PC E2E 테스트
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(slide, "04-B. PC E2E 테스트 (pc_e2e_test_v2.spec.ts)", "PC E2E — 18 Test Groups")

# 기본 정보
info_items = [
    ("baseURL", "https://store.hanssem.com"),
    ("샘플 상품", "#837513"),
    ("결과 파일", "public/pc_e2e.json"),
    ("ID", "pc-e2e"),
    ("cases", "pass + fail 전체 포함"),
]
for i, (k, v) in enumerate(info_items):
    x = 0.3 + (i%3)*4.3
    y = 1.3 if i<3 else 1.85
    add_rect(slide, x, y, 4.1, 0.45, fill=LIGHT_BLUE, line=HANSSEM_BLUE, line_w=0.3)
    add_text_box(slide, f"{k}: {v}", x+0.1, y+0.05, 3.9, 0.35, font_size=10, bold=False, color=DARK_TEXT)

# 18개 그룹 표
headers = ["#", "그룹명", "주요 검증 항목"]
rows = [
    ("1",  "메인 페이지",       "로딩/타이틀, 헤더(GNB), 푸터(PC)"),
    ("2",  "카테고리 네비게이션", "가구/홈리빙·인테리어 진입, GNB 링크 클릭(PC)"),
    ("3",  "상품 목록",         "상품 링크 노출, 상품 클릭→상세 이동"),
    ("4",  "상품 상세",         "상품명(H1), 가격, 구매/장바구니 버튼(PC), 이미지"),
    ("5",  "검색",              "검색아이콘→입력→결과 이동(PC), 영문 키워드"),
    ("6",  "로그인",            "GNB 로그인 클릭→mall.hanssem.com(PC), 로그인 폼"),
    ("7",  "장바구니",          "GNB 장바구니 클릭(PC), 직접 URL redirect"),
    ("8",  "HTTP 응답",         "메인/가구/인테리어/검색/상품상세/매장 6개 200 응답"),
    ("9",  "매장 찾기",         "직접 진입, GNB 매장찾기 클릭(PC)"),
    ("10", "상품 탭 전환",      "후기·문의·배송 탭 전환"),
    ("11", "정렬/필터",         "모든필터(PC), 낮은가격순 정렬(PC), 0건 검색"),
    ("12", "인테리어 서브",     "카테고리, 시공사례, 무료견적상담, 기획전 200"),
    ("13", "예외 페이지",       "없는 상품 redirect (5xx 미발생 확인)"),
    ("14", "옵션 선택 레이어",  "옵션 버튼 노출(PC), 드롭다운 노출(PC)"),
    ("15", "시공사례 상세",     "목록 노출, 첫 번째 항목 클릭→상세 진입"),
    ("16", "전문가 찾기",       "페이지 진입, 전문가 카드 목록 노출"),
    ("17", "매장 검색/상세",    "검색 입력 필드(PC), 지역 필터/매장 목록 노출"),
    ("18", "붙박이장 셀프플래너","플래너 링크, 붙박이장 카테고리 상품 노출"),
]
add_table(slide, headers, rows, 0.3, 2.5, 12.73, 4.75, col_widths=[0.5, 3.0, 6.0])

# ══════════════════════════════════════════════════════════════
# 슬라이드 8: MW E2E 테스트
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(slide, "04-C. MW E2E 테스트 (mw_e2e_test.spec.ts)", "MW E2E — 13 Test Groups")

# 기본 정보
add_rect(slide, 0.3, 1.25, 8.5, 0.85, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_multiline(slide, [
    "baseURL: https://m.store.hanssem.com  |  샘플 상품: #837513  |  결과 파일: public/mw_e2e.json",
    "cases: results.json에 fail만, mw_e2e.json에 pass+fail 전체  |  hasPublished 플래그로 afterAll 중복 방지",
], 0.45, 1.3, 8.2, 0.75, font_size=10, color=DARK_TEXT)

# 구조 주의사항
add_rect(slide, 9.0, 1.25, 4.0, 0.85, fill=RGBColor(0xFF, 0xF3, 0xCD), line=ACCENT_ORANGE, line_w=0.5)
add_text_box(slide, "⚠️ afterAll 래퍼 필수", 9.1, 1.28, 3.8, 0.35, font_size=10, bold=True, color=ACCENT_ORANGE)
add_text_box(slide, "전체를 test.describe(\"MW E2E 테스트\") 하나로 감싸야 함\nafterAll 중복 실행 방지 (결과 덮어씌움 버그)", 9.1, 1.6, 3.8, 0.45, font_size=9, color=DARK_TEXT)

# 13개 그룹 표
headers = ["#", "그룹명", "주요 검증 항목", "PC와 차이점"]
rows = [
    ("1",  "메인 페이지",      "로딩/타이틀, 헤더",                    "푸터 없음 (BNB 사용)"),
    ("2",  "카테고리 네비게이션", "가구/홈리빙·인테리어 진입",           "GNB 링크 클릭 없음"),
    ("3",  "상품 목록",        "상품 링크 노출, 상품 클릭→상세 이동",   "동일"),
    ("4",  "상품 상세",        "상품명(H1), 가격, 이미지",              "구매버튼 검증 없음 (MW UI 다름)"),
    ("5",  "검색",             "검색 결과 로딩, 영문 키워드",           "검색아이콘 GNB 클릭 없음"),
    ("6",  "로그인",           "로그인 폼 노출",                        "returnUrl = m.store.hanssem.com"),
    ("7",  "장바구니",         "직접 URL redirect",                    "mall.hanssem.com/m/morder 경로"),
    ("8",  "HTTP 응답",        "6개 주요 페이지 200 응답",              "동일"),
    ("9",  "매장 찾기",        "직접 진입 및 타이틀",                   "GNB 클릭 없음"),
    ("10", "상품 탭 전환",     "후기 탭 전환",                          "문의/배송 탭 없음 (MW UI 다름)"),
    ("11", "검색 결과",        "0건 검색결과 확인",                     "동일"),
    ("12", "인테리어 서브",    "카테고리(MW URL), 무료견적상담",        "m.store.hanssem.com URL 사용"),
    ("13", "예외 페이지",      "없는 상품 redirect",                    "동일"),
]
add_table(slide, headers, rows, 0.3, 2.25, 12.73, 4.95, col_widths=[0.5, 2.5, 4.5, 3.0])

# ══════════════════════════════════════════════════════════════
# 슬라이드 9: 결과 데이터 & 대시보드
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(slide, "05. 결과 데이터 & 대시보드", "Results Data & Dashboard")

# results.json 구조
add_rect(slide, 0.3, 1.25, 6.0, 5.9, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 0.3, 1.25, 6.0, 0.38, fill=HANSSEM_BLUE)
add_text_box(slide, "results.json 구조 (대시보드 단일 소스)", 0.45, 1.28, 5.7, 0.32, font_size=11, bold=True, color=WHITE)
json_lines = [
    '{',
    '  "lastUpdated": "2026. 4. 20. ...",',
    '  "codeLastModified": "2026. 04. 20. 09:30",',
    '  "incidents": {',
    '    "total": 30, "hotfix": 30,',
    '    "feCount": 14, "beCount": 15, "appCount": 1,',
    '    "byMonth": [ { "month": "2026-04",',
    '      "count": 11, "fe": 6, "be": 5,',
    '      "list": [ { "no": 11, "date": "...",',
    '        "category": "FE", "assignee": "...",',
    '        "summary": "...", "jiraLink": "..." } ] } ]',
    '  },',
    '  "reports": [ {',
    '    "id": "pc-landing",',
    '    "title": "운영환경 PC 500개 랜딩 테스트",',
    '    "total": 500, "pass": 498, "fail": 2,',
    '    "passRate": "99.6",',
    '    "cases": [ /* fail만 */ ]',
    '  } ]',
    '}',
]
add_multiline(slide, json_lines, 0.42, 1.7, 5.7, 5.3, font_size=9, color=RGBColor(0x1A, 0x56, 0x1A))

# 파일 매핑
add_rect(slide, 6.55, 1.25, 6.45, 2.1, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 6.55, 1.25, 6.45, 0.38, fill=RGBColor(0x00, 0x80, 0xC0))
add_text_box(slide, "파일 매핑 및 cases 포함 범위", 6.7, 1.28, 6.1, 0.32, font_size=11, bold=True, color=WHITE)
f_headers = ["ID", "상세 파일", "results.json", "상세 화면"]
f_rows = [
    ("pc-landing", "pc_500.json", "fail만", "전체"),
    ("pc-e2e",     "pc_e2e.json", "전체",  "전체"),
    ("mw-landing", "mw_500.json", "fail만", "전체"),
    ("mw-e2e",     "mw_e2e.json", "fail만", "전체"),
]
add_table(slide, f_headers, f_rows, 6.55, 1.68, 6.45, 1.62, col_widths=[2.0, 2.5, 1.5, 1.5])

# 대시보드 기능
add_rect(slide, 6.55, 3.5, 6.45, 3.65, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 6.55, 3.5, 6.45, 0.38, fill=HANSSEM_BLUE)
add_text_box(slide, "대시보드 (index.html) 주요 기능", 6.7, 3.53, 6.1, 0.32, font_size=11, bold=True, color=WHITE)
add_multiline(slide, [
    "📊 테스트 결과 카드 (PC 랜딩 / PC E2E / MW 랜딩 / MW E2E)",
    "📈 월별 긴급배포 막대 그래프 (FE=파랑, BE=주황, APP=초록)",
    "🔴 fail > 0이면 '실패 감지' 배지 표시",
    "⏱️  실행 중 로딩 UI:",
    "   · 랜딩: count/total × 100 실제 % 표시 (15초 갱신)",
    "   · E2E: 스피너만 표시 (% 없음)",
    "   · 완료 카드: '완료: HH:MM' 표시",
    "   · 종료 감지 → 5초 후 결과 자동 로드",
    "",
    "폴링 타이머",
    "  progressTimer: 30초  GitHub 워크플로우 상태 확인",
    "  phaseTimer:    15초  progress.json 폴링 (단계 감지)",
    "  animTimer:     10초  % 수치 갱신",
], 6.7, 3.95, 6.2, 3.15, font_size=10, color=DARK_TEXT)

# ══════════════════════════════════════════════════════════════
# 슬라이드 10: GitHub Actions 워크플로우
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(slide, "06. GitHub Actions 워크플로우", "CI/CD Workflow")

# daily-trigger
add_rect(slide, 0.3, 1.25, 5.9, 2.5, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 0.3, 1.25, 5.9, 0.38, fill=ACCENT_ORANGE)
add_text_box(slide, "daily-trigger.yml", 0.45, 1.28, 5.6, 0.32, font_size=11, bold=True, color=WHITE)
add_multiline(slide, [
    "스케줄:  0 23 * * * (UTC) = KST 08:00",
    "방식:    PAT_TOKEN으로 GitHub API 호출",
    "         → playwright.yml workflow_dispatch 트리거",
    "Secret:  PAT_TOKEN",
    "",
    "분리 이유:",
    "  GitHub 내장 schedule cron의 신뢰도 문제",
    "  (지연/미실행 사례 발생) → 별도 트리거로 분리",
], 0.45, 1.7, 5.6, 2.0, font_size=10, color=DARK_TEXT)

# playwright.yml 단계
add_rect(slide, 6.45, 1.25, 6.55, 2.5, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 6.45, 1.25, 6.55, 0.38, fill=HANSSEM_BLUE)
add_text_box(slide, "playwright.yml — 환경 설정", 6.6, 1.28, 6.2, 0.32, font_size=11, bold=True, color=WHITE)
add_multiline(slide, [
    "트리거:   workflow_dispatch (수동 또는 daily-trigger)",
    "타임아웃: 300분",
    "Secrets:  GITHUB_TOKEN, PAT_TOKEN",
    "",
    "환경변수:",
    "  CI: ''   → 잔디 알림 임시 비활성화",
    "             (삭제하면 알림 재활성화)",
    "  GITHUB_TOKEN → publishResults, updateProgress에 사용",
], 6.6, 1.7, 6.2, 2.0, font_size=10, color=DARK_TEXT)

# 실행 단계 상세
steps_detail = [
    ("1단계", "codeLastModified 기록", HANSSEM_BLUE,
     "최근 non-auto 커밋 날짜를 results.json에 기록\n(auto update 커밋 제외하여 실제 코드 변경일 표시)"),
    ("2단계", "의존성 설치", HANSSEM_BLUE,
     "npm ci\nnpx playwright install --with-deps"),
    ("3단계", "테스트 실행", GREEN,
     "npx playwright test\n(continue-on-error: true — 실패해도 이후 단계 계속)"),
    ("4단계", "결과 커밋", HANSSEM_BLUE,
     "git pull --rebase origin main && git push\n(force push 금지 — Actions 권한 오류 발생)"),
    ("5단계", "Pages 배포", HANSSEM_BLUE,
     "peaceiris/actions-gh-pages@v3\npublic/ → gh-pages 브랜치"),
    ("6단계", "아티팩트 업로드", MID_GRAY,
     "playwright-report, fail_evidence\n(30일 보관)"),
]

for i, (n, t, color, desc) in enumerate(steps_detail):
    col = i % 3
    row = i // 3
    x = 0.3 + col * 4.3
    y = 3.95 + row * 1.75
    add_rect(slide, x, y, 4.1, 1.6, fill=WHITE, line=color, line_w=0.5)
    add_rect(slide, x, y, 4.1, 0.38, fill=color)
    add_text_box(slide, f"{n}  {t}", x+0.1, y+0.04, 3.9, 0.3, font_size=10, bold=True, color=WHITE)
    add_text_box(slide, desc, x+0.1, y+0.43, 3.9, 1.1, font_size=9.5, color=DARK_TEXT)

# ══════════════════════════════════════════════════════════════
# 슬라이드 11: 기능 커버리지 분석
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(slide, "07. 기능 커버리지 분석", "Feature Coverage Analysis — 66 items")

# 요약 수치 카드
kpis2 = [
    ("66개", "전체 기본기능 항목", HANSSEM_BLUE),
    ("8개",  "완전 커버",          GREEN),
    ("11개", "부분 커버",          YELLOW),
    ("47개", "미커버",             RED),
    ("약 21%", "총 커버리지",      HANSSEM_BLUE),
]
for i, (val, lbl, color) in enumerate(kpis2):
    x = 0.3 + i * 2.55
    add_rect(slide, x, 1.25, 2.35, 1.2, fill=WHITE, line=color, line_w=0.5)
    add_rect(slide, x, 1.25, 2.35, 0.08, fill=color)
    add_text_box(slide, val, x, 1.42, 2.35, 0.6,
                 font_size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, lbl, x, 2.02, 2.35, 0.38,
                 font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)

# 영역별 커버리지 표
headers = ["기능 영역", "전체", "완전", "부분", "미커버", "커버율"]
rows_cov = [
    ("회원/인증",   "9",  "0", "1",  "8",  "약 6%"),
    ("주문/결제",  "13",  "0", "0", "13",  "0%"),
    ("메인/네비게이션", "11", "2", "4",  "5", "약 36%"),
    ("상품",       "16",  "3", "2", "11",  "약 25%"),
    ("인테리어",    "9",  "2", "3",  "4",  "약 39%"),
    ("마이페이지",  "8",  "0", "0",  "8",  "0%"),
    ("합계",       "66",  "7","10", "49",  "약 21%"),
]
add_table(slide, headers, rows_cov, 0.3, 2.65, 6.3, 2.9, col_widths=[2.5, 0.8, 0.8, 0.8, 0.8, 1.0])

# 미커버 원인 분석
add_rect(slide, 6.85, 2.65, 6.15, 2.9, fill=WHITE, line=RED, line_w=0.3)
add_rect(slide, 6.85, 2.65, 6.15, 0.38, fill=RED)
add_text_box(slide, "미커버 주요 원인", 7.0, 2.68, 5.9, 0.32, font_size=11, bold=True, color=WHITE)
cause_items = [
    ("로그인 세션 필요",   "약 25개", "장바구니·주문·결제·마이페이지 전체"),
    ("실결제 필요",       "13개",    "신용카드·페이코·카카오페이 등"),
    ("외부 OAuth 필요",   "5개",     "카카오/네이버/페이코/애플 간편가입"),
    ("특수 상품 유형",    "4개",     "커스텀·커튼/블라인드·패키지"),
    ("미구현 시나리오",   "약 5개",  "매거진·사이드메뉴·RD상세"),
]
for i, (cause, cnt, note) in enumerate(cause_items):
    y = 3.1 + i * 0.43
    add_rect(slide, 6.85, y, 0.9, 0.35, fill=RED)
    add_text_box(slide, cnt, 6.85, y+0.04, 0.9, 0.27, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, f"{cause} — {note}", 7.82, y+0.04, 5.1, 0.27, font_size=9.5, color=DARK_TEXT)

# 커버리지 확대 방향
add_rect(slide, 0.3, 5.7, 12.73, 1.55, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 0.3, 5.7, 12.73, 0.38, fill=HANSSEM_BLUE)
add_text_box(slide, "커버리지 확대 방향 (우선순위 순)", 0.45, 5.73, 12.4, 0.32, font_size=11, bold=True, color=WHITE)
priorities = [
    ("높음", "장바구니/구매 플로우",      "테스트 전용 계정 세션 저장 후 재사용"),
    ("높음", "사이드(햄버거) 메뉴 이동",  "MW E2E 시나리오 추가"),
    ("중간", "검색 내 상세 이동 (4종)",   "검색 결과 클릭 시나리오 추가"),
    ("중간", "매거진, RD 상세",           "E2E 시나리오 추가"),
    ("낮음", "결제/주문 완료",            "테스트 환경 결제 수단 설정 필요"),
]
for i, (p, feat, how) in enumerate(priorities):
    x = 0.5 + i * 2.5
    color = RED if p == "높음" else (YELLOW if p == "중간" else GREEN)
    add_text_box(slide, p, x, 6.13, 0.6, 0.28, font_size=9, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, feat, x+0.65, 6.12, 1.75, 0.3, font_size=9, bold=True, color=DARK_TEXT)
    add_text_box(slide, how, x+0.65, 6.41, 1.75, 0.3, font_size=8.5, color=MID_GRAY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 12: 운영 가이드
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(slide, "08. 운영 가이드", "Operations Guide")

# 장애 리포트 추가
add_rect(slide, 0.3, 1.25, 6.0, 2.8, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 0.3, 1.25, 6.0, 0.38, fill=HANSSEM_BLUE)
add_text_box(slide, "장애 리포트 추가 방법", 0.45, 1.28, 5.7, 0.32, font_size=11, bold=True, color=WHITE)
add_multiline(slide, [
    "핫픽스/장애 발생 시 public/results.json 직접 수정",
    "",
    "1. incidents.byMonth[해당월].list에 항목 추가",
    "   (no: 월별 순번, 최신이 가장 위)",
    "2. 해당 월 count / fe / be / app 업데이트",
    "3. 최상위 total / feCount / beCount / appCount 업데이트",
    "4. lastUpdated / period 날짜 갱신",
    "5. git commit + push",
], 0.45, 1.7, 5.7, 2.25, font_size=10, color=DARK_TEXT)

# 잔디 알림
add_rect(slide, 6.55, 1.25, 6.45, 2.8, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 6.55, 1.25, 6.45, 0.38, fill=RGBColor(0x00, 0x80, 0xC0))
add_text_box(slide, "잔디(Jandi) 알림", 6.7, 1.28, 6.1, 0.32, font_size=11, bold=True, color=WHITE)
add_multiline(slide, [
    "전송 조건: process.env.CI가 truthy일 때만",
    "전송 시점: 각 테스트 afterAll 완료 후 (4회)",
    "현재 상태: CI='' 설정으로 임시 비활성화",
    "",
    "재활성화 방법:",
    "  playwright.yml의 아래 블록을 삭제",
    "  env:",
    "    CI: ''",
    "",
    "알림 내용: 통과율, 총 건수,",
    "           실패 목록(최대 10개), 대시보드 링크",
], 6.7, 1.7, 6.1, 2.25, font_size=10, color=DARK_TEXT)

# PAT 토큰
add_rect(slide, 0.3, 4.2, 6.0, 2.0, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 0.3, 4.2, 6.0, 0.38, fill=ACCENT_ORANGE)
add_text_box(slide, "GitHub PAT 토큰 관리", 0.45, 4.23, 5.7, 0.32, font_size=11, bold=True, color=WHITE)
add_multiline(slide, [
    "대시보드 index.html에 하드코딩",
    "(split array로 GitHub push protection 우회)",
    "",
    'const GH_TOKEN = ["ghp","_0yXkI87t","kJ1wcUmh",',
    '                  "lLIfYFs2","m4OH8243","MWD7"].join("");',
    "",
    "토큰 만료 시: 새 토큰 발급 → 위 배열 업데이트",
], 0.45, 4.65, 5.7, 1.48, font_size=10, color=DARK_TEXT)

# Git 배포 규칙
add_rect(slide, 6.55, 4.2, 6.45, 2.0, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_rect(slide, 6.55, 4.2, 6.45, 0.38, fill=HANSSEM_BLUE)
add_text_box(slide, "Git 배포 규칙", 6.7, 4.23, 6.1, 0.32, font_size=11, bold=True, color=WHITE)
add_multiline(slide, [
    "❌ git push --force 금지",
    "   (GitHub Actions 권한 오류 발생)",
    "✅ git pull --rebase origin main && git push 사용",
    "❌ 테스트 실행 중 중간 푸시 금지 (Actions 충돌)",
    "✅ progress.json은 Contents API로 gh-pages에 직접 기록",
    "   (git push 아님)",
], 6.7, 4.65, 6.1, 1.48, font_size=10, color=DARK_TEXT)

# 패키지
add_rect(slide, 0.3, 6.32, 12.73, 0.95, fill=WHITE, line=HANSSEM_BLUE, line_w=0.3)
add_text_box(slide, "패키지:", 0.45, 6.38, 1.5, 0.35, font_size=10, bold=True, color=HANSSEM_BLUE)
add_text_box(slide,
    "@playwright/test ^1.59.1  |  @types/node ^20.19.39  |  axios ^1.6.0  "
    "|  tsconfig.json 불필요 (Playwright가 TypeScript 직접 처리)",
    1.9, 6.38, 11.0, 0.35, font_size=10, color=DARK_TEXT)
add_text_box(slide, "대시보드:", 0.45, 6.72, 1.5, 0.35, font_size=10, bold=True, color=HANSSEM_BLUE)
add_text_box(slide,
    "https://daewon82.github.io/hanssem-qa-system/",
    1.9, 6.72, 11.0, 0.35, font_size=10, color=HANSSEM_BLUE)

prs.save("/Users/dw/hanssem-qa-system/한샘몰_QA시스템_기술문서.pptx")
print("✅ PPT 저장 완료: 한샘몰_QA시스템_기술문서.pptx")
